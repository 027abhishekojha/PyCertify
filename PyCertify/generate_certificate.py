import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas
from io import BytesIO

CHARACTER_SPACING = 0.25  # Adjust this value as needed

def replace_text_in_shape(shape, old_text, new_text):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = run.text.replace(old_text, new_text)

def calculate_run_position(shape, run):
    if run.font.size is not None:
        left = Inches(shape.left) + Pt(run.font.size) * CHARACTER_SPACING
        top = Inches(shape.top)
        return left, top
    return None, None

def generate_certificates(template_path, excel_path, output_folder, num_copies, pptx_output_folder, pdf_output_folder):
    # Load Excel data
    df = pd.read_excel(excel_path)

    # Create folders for PowerPoint and PDF files
    os.makedirs(pptx_output_folder, exist_ok=True)
    os.makedirs(pdf_output_folder, exist_ok=True)

    for index, row in df.iterrows():
        # Load PowerPoint template
        prs = Presentation(template_path)

        # Find and replace placeholder text
        placeholder_text = '<<NAME_PLACEHOLDER>>'
        for slide in prs.slides:
            for shape in slide.shapes:
                replace_text_in_shape(shape, placeholder_text, str(row['ParticipantName']))

        # Save the new presentation
        pptx_output_path = os.path.join(pptx_output_folder, f"Certificate_{index + 1}.pptx")
        prs.save(pptx_output_path)

        # Create PDF from PPTX using reportlab
        pdf_output_path = os.path.join(pdf_output_folder, f"Certificate_{index + 1}.pdf")
        pdf_stream = BytesIO()
        canvas_obj = canvas.Canvas(pdf_stream)

        # Extract text content from PowerPoint and add it to PDF
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            left, top = calculate_run_position(shape, run)
                            if left is not None and top is not None:
                                canvas_obj.drawString(left, top, run.text)

        canvas_obj.save()

        pdf_stream.seek(0)
        with open(pdf_output_path, 'wb') as pdf_file:
            pdf_file.write(pdf_stream.read())

    print(f"{num_copies} certificates generated successfully.")

if __name__ == "__main__":
    # Set your template path, Excel path, output folders, and the number of copies
    template_path = 'C:/Users/aayus/OneDrive/Desktop/New folder/PyCertify/Certificate_template.pptx'
    excel_path = 'C:/Users/aayus/OneDrive/Desktop/New folder/PyCertify/participants.xlsx'
    output_folder = 'C:/Users/aayus/OneDrive/Desktop/New folder/PyCertify'
    pptx_output_folder = os.path.join(output_folder, 'pptx_output')
    pdf_output_folder = os.path.join(output_folder, 'pdf_output')

    num_copies = 100  # Adjust as needed

    generate_certificates(template_path, excel_path, output_folder, num_copies, pptx_output_folder, pdf_output_folder)
